import streamlit as st
from openai import OpenAI
from pptx import Presentation
from PIL import Image, ImageDraw
import io
import base64
import fitz  # PyMuPDF

# === CONFIGURACIÓN ===
st.set_page_config(page_title="Validación AI - PPTX", layout="wide")
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

st.title("📊 Revisión Automática de Reportes CFO (GPT-4o + Manual)")

# === SUBIDA DE ARCHIVOS ===
pptx_file = st.file_uploader("📄 Sube tu reporte PowerPoint (.pptx)", type=["pptx"])

# === FUNCIONES ===

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return "\n".join([page.get_text() for page in doc])

def render_slide_as_image(slide, width=1280, height=720):
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    y_offset = 40
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            draw.text((50, y_offset), shape.text, fill="black")
            y_offset += 40
    return img

def pptx_to_base64_images(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    base64_images = []
    for slide in prs.slides:
        img = render_slide_as_image(slide)
        with io.BytesIO() as output:
            img.save(output, format="PNG")
            b64 = base64.b64encode(output.getvalue()).decode("utf-8")
            base64_images.append(b64)
    return base64_images

def create_gpt4o_payload(images_b64, max_slides=40):
    inputs = []
    for b64 in images_b64[:max_slides]:
        inputs.append({
            "type": "image_url",
            "image_url": {
                "url": f"data:image/png;base64,{b64}"
            }
        })
    return inputs

# === CARGAR MANUAL PDF Y PROCESAR ===
manual_text = extract_text_from_pdf("Monthly. Quality checks.pdf")
manual_text_short = manual_text[:15000]  # evitar exceder tokens

# === PROCESAMIENTO DEL REPORTE ===
if pptx_file:
    pptx_bytes = pptx_file.read()
    st.info("🔄 Convirtiendo diapositivas en imágenes...")
    slides_b64 = pptx_to_base64_images(pptx_bytes)
    slide_inputs = create_gpt4o_payload(slides_b64, max_slides=40)

    st.success(f"✅ {len(slide_inputs)} diapositivas procesadas.")

    st.info("📨 Enviando a GPT-4o para análisis...")

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Tengo el siguiente reporte (formato PowerPoint) y el manual de control de calidad de Monthly. "
                        "Evalúa todas las diapositivas conforme a dicho manual. "
                        "Identifica errores, inconsistencias y oportunidades de mejora por sección: Resumen, Gráficos, Tablas, Notación, "
                        "Métricas, Comparativos, Ciclo de Conversión, etc. Termina con una tabla o lista de recomendaciones para el equipo. \n\n"
                        "Aquí está el contenido del manual:\n\n"
                        f"{manual_text_short}"
                    )
                },
                {
                    "role": "user",
                    "content": [
                        { "type": "text", "text": "Este es el reporte a revisar, generado automáticamente:" },
                        *slide_inputs
                    ]
                }
            ],
            temperature=0.3,
            max_tokens=3000
        )

        resultado = response.choices[0].message.content

        st.markdown("### ✅ Evaluación del Reporte:")
        st.write(resultado)

        st.download_button(
            label="📥 Descargar evaluación como .txt",
            data=resultado,
            file_name="evaluacion_reporte.txt",
            mime="text/plain"
        )

    except Exception as e:
        st.error("❌ Ocurrió un error al llamar a OpenAI.")
        st.exception(e)
